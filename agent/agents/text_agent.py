"""
Text Agent - Specialized processor for text-based documents
Handles: TXT, DOCX, DOC, PPTX, PPT, RTF, ODT files
Enhanced with advanced structure analysis and key concept extraction
"""
import os
import json
import re
from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass, asdict
import docx
from pptx import Presentation
import boto3

try:
    from ..config.model_manager import model_config_manager
except ImportError:
    # Fallback for when not running as package
    import sys
    import os
    sys.path.append(os.path.dirname(os.path.dirname(__file__)))
    from config.model_manager import model_config_manager


@dataclass
class TextStructure:
    """Represents the structure of extracted text content."""
    headings: List[Dict[str, Any]]
    paragraphs: List[str]
    lists: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    formatting: Dict[str, Any]


@dataclass
class KeyConcept:
    """Represents a key concept with confidence scoring."""
    concept: str
    definition: str
    importance: float  # 0.0 to 1.0
    confidence: float  # 0.0 to 1.0
    context: str


@dataclass
class EnhancedTextAnalysis:
    """Enhanced analysis results for text content."""
    structure: TextStructure
    key_concepts: List[KeyConcept]
    educational_metadata: Dict[str, Any]
    confidence_scores: Dict[str, float]
    processing_metrics: Dict[str, Any]


class TextAgent:
    """Enhanced specialized agent for processing text documents."""
    
    def __init__(self, region: str = "us-east-1"):
        self.region = region
        self.bedrock_client = boto3.client('bedrock-runtime', region_name=region)
        self.supported_extensions = ['.txt', '.docx', '.doc', '.pptx', '.ppt', '.rtf', '.odt']
        self.model_manager = model_config_manager
    
    def can_process(self, file_path: str) -> bool:
        """Check if this agent can process the given file."""
        return any(file_path.lower().endswith(ext) for ext in self.supported_extensions)
    
    def _resolve_file_path(self, file_path: str) -> str:
        """Resolve file path relative to project structure."""
        from pathlib import Path
        
        file_path_obj = Path(file_path)
        print(f"🔍 TEXT Agent - Resolving file path: {file_path_obj}")
        
        # If the path exists as-is, use it
        if file_path_obj.exists():
            print(f"✅ TEXT Agent - Found file at original path: {file_path_obj}")
            return str(file_path_obj)
        
        # Try in backend directory (most common case)
        backend_path = Path("backend") / file_path_obj
        if backend_path.exists():
            print(f"✅ TEXT Agent - Found file at backend path: {backend_path}")
            return str(backend_path)
        
        # Try relative to backend directory (from agent directory)
        backend_relative_path = Path("../backend") / file_path_obj
        if backend_relative_path.exists():
            print(f"✅ TEXT Agent - Found file at backend relative path: {backend_relative_path}")
            return str(backend_relative_path)
        
        # Return original path if nothing works
        print(f"❌ TEXT Agent - Could not resolve file path, using original: {file_path_obj}")
        return file_path
    
    async def process_file(self, file_path: str, user_id: str) -> Dict[str, Any]:
        """Process a text document file with enhanced extraction."""
        try:
            print(f"📄 Enhanced TEXT Agent processing: {file_path}")

            # Resolve file path
            resolved_path = self._resolve_file_path(file_path)

            # Extract structured content based on file type
            structured_content = await self._extract_structured_content(resolved_path)

            # Perform enhanced AI analysis
            enhanced_analysis = await self._perform_enhanced_analysis(
                structured_content, resolved_path
            )
            
            # Prepare enhanced result
            result = {
                "agent_type": "text",
                "file_path": file_path,
                "status": "completed",
                "content": {
                    "text": self._create_smart_preview(structured_content.get("raw_text", ""), 8000),
                    "full_text": structured_content.get("raw_text", ""),
                    "word_count": len(structured_content.get("raw_text", "").split()),
                    "char_count": len(structured_content.get("raw_text", "")),
                    "structure": structured_content.get("structure"),
                    "formatting": structured_content.get("formatting", {})
                },
                "enhanced_analysis": asdict(enhanced_analysis) if enhanced_analysis else {},
                "metadata": {
                    "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
                    "file_type": Path(file_path).suffix.lower(),
                    "processed_by": "enhanced_text_agent",
                    "model_used": self._get_current_model_id(),
                    "processing_version": "2.0"
                }
            }
            
            print(f"✅ Enhanced TEXT Agent completed: {file_path}")
            return result
            
        except Exception as e:
            print(f"❌ Enhanced TEXT Agent error processing {file_path}: {e}")
            # Try fallback processing
            fallback_result = await self._fallback_processing(file_path, str(e))
            return fallback_result
    
    async def _extract_structured_content(self, file_path: str) -> Dict[str, Any]:
        """Extract structured content with advanced text analysis."""
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.txt':
                return await self._extract_txt_structure(file_path)
            elif file_ext in ['.docx', '.doc']:
                return await self._extract_docx_structure(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return await self._extract_pptx_structure(file_path)
            else:
                # Fallback: try to read as plain text
                return await self._extract_txt_structure(file_path)
                
        except Exception as e:
            print(f"⚠️ Error extracting structured content from {file_path}: {e}")
            return {
                "raw_text": f"Error extracting content: {str(e)}",
                "structure": TextStructure([], [], [], [], {}),
                "formatting": {},
                "extraction_error": str(e)
            }
    
    async def _extract_txt_structure(self, file_path: str) -> Dict[str, Any]:
        """Extract structure from plain text files."""
        try:
            # Resolve file path first
            resolved_path = self._resolve_file_path(file_path)
            print(f"📝 TEXT Agent processing resolved path: {resolved_path}")
            
            with open(resolved_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Analyze text structure
            structure = self._analyze_text_structure(content)
            
            return {
                "raw_text": content,
                "structure": structure,
                "formatting": {"type": "plain_text"},
                "extraction_method": "text_analysis"
            }
            
        except Exception as e:
            raise Exception(f"Failed to extract text structure: {e}")
    
    async def _extract_docx_structure(self, file_path: str) -> Dict[str, Any]:
        """Extract structure from DOCX files with formatting preservation."""
        try:
            doc = docx.Document(file_path)
            
            headings = []
            paragraphs = []
            lists = []
            tables = []
            formatting_info = {"styles": [], "fonts": set()}
            
            # Extract paragraphs with style information
            for para in doc.paragraphs:
                if para.text.strip():
                    # Check if it's a heading
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1
                        headings.append({
                            "text": para.text.strip(),
                            "level": level,
                            "style": para.style.name
                        })
                    else:
                        paragraphs.append(para.text.strip())
                    
                    # Track formatting
                    formatting_info["styles"].append(para.style.name)
                    for run in para.runs:
                        if run.font.name:
                            formatting_info["fonts"].add(run.font.name)
            
            # Extract tables
            for table in doc.tables:
                table_data = {
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "content": []
                }
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data["content"].append(row_data)
                tables.append(table_data)
            
            # Convert sets to lists for JSON serialization
            formatting_info["fonts"] = list(formatting_info["fonts"])
            
            raw_text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
            
            structure = TextStructure(
                headings=headings,
                paragraphs=paragraphs,
                lists=lists,  # Will be enhanced in future iterations
                tables=tables,
                formatting=formatting_info
            )
            
            return {
                "raw_text": raw_text,
                "structure": structure,
                "formatting": formatting_info,
                "extraction_method": "docx_analysis"
            }
            
        except Exception as e:
            raise Exception(f"Failed to extract DOCX structure: {e}")
    
    async def _extract_pptx_structure(self, file_path: str) -> Dict[str, Any]:
        """Extract structure from PowerPoint files."""
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            all_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_data = {
                    "slide_number": i + 1,
                    "title": "",
                    "content": [],
                    "notes": ""
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Try to identify title vs content
                        if shape.placeholder_format and shape.placeholder_format.type == 1:  # Title placeholder
                            slide_data["title"] = shape.text.strip()
                        else:
                            slide_data["content"].append(shape.text.strip())
                        
                        all_text.append(shape.text.strip())
                
                # Extract slide notes if available
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
                
                slides_content.append(slide_data)
            
            raw_text = '\n'.join(all_text)
            
            # Create structure representation
            headings = [{"text": slide["title"], "level": 1, "slide": slide["slide_number"]} 
                       for slide in slides_content if slide["title"]]
            
            paragraphs = []
            for slide in slides_content:
                paragraphs.extend(slide["content"])
            
            structure = TextStructure(
                headings=headings,
                paragraphs=paragraphs,
                lists=[],  # Will be enhanced in future iterations
                tables=[],
                formatting={"type": "presentation", "slides": len(slides_content)}
            )
            
            return {
                "raw_text": raw_text,
                "structure": structure,
                "slides": slides_content,
                "formatting": {"type": "presentation", "total_slides": len(slides_content)},
                "extraction_method": "pptx_analysis"
            }
            
        except Exception as e:
            raise Exception(f"Failed to extract PPTX structure: {e}")
    
    def _analyze_text_structure(self, text: str) -> TextStructure:
        """Analyze structure of plain text content."""
        lines = text.split('\n')
        
        headings = []
        paragraphs = []
        lists = []
        current_paragraph = []
        
        for line in lines:
            line = line.strip()
            if not line:
                if current_paragraph:
                    paragraphs.append(' '.join(current_paragraph))
                    current_paragraph = []
                continue
            
            # Detect headings (lines that are short, capitalized, or have special formatting)
            if self._is_likely_heading(line):
                if current_paragraph:
                    paragraphs.append(' '.join(current_paragraph))
                    current_paragraph = []
                
                headings.append({
                    "text": line,
                    "level": self._estimate_heading_level(line),
                    "style": "detected"
                })
            
            # Detect list items
            elif self._is_list_item(line):
                if current_paragraph:
                    paragraphs.append(' '.join(current_paragraph))
                    current_paragraph = []
                
                lists.append({
                    "text": line,
                    "type": "bullet" if line.startswith(('•', '-', '*')) else "numbered"
                })
            
            else:
                current_paragraph.append(line)
        
        # Add final paragraph if exists
        if current_paragraph:
            paragraphs.append(' '.join(current_paragraph))
        
        return TextStructure(
            headings=headings,
            paragraphs=paragraphs,
            lists=lists,
            tables=[],  # Plain text doesn't have tables
            formatting={"type": "plain_text"}
        )
    
    def _is_likely_heading(self, line: str) -> bool:
        """Determine if a line is likely a heading."""
        # Check various heading indicators
        if len(line) < 5:  # Very short lines might be headings
            return False
        
        if len(line) > 100:  # Very long lines are unlikely to be headings
            return False
        
        # Check for common heading patterns
        heading_patterns = [
            r'^[A-Z][A-Z\s]+$',  # ALL CAPS
            r'^\d+\.\s+[A-Z]',   # Numbered headings
            r'^[A-Z][a-z]+(\s[A-Z][a-z]+)*$',  # Title Case
            r'^#+\s+',           # Markdown-style headings
        ]
        
        for pattern in heading_patterns:
            if re.match(pattern, line):
                return True
        
        return False
    
    def _estimate_heading_level(self, line: str) -> int:
        """Estimate the heading level based on formatting."""
        if line.startswith('###'):
            return 3
        elif line.startswith('##'):
            return 2
        elif line.startswith('#'):
            return 1
        elif re.match(r'^\d+\.\d+\.\d+', line):
            return 3
        elif re.match(r'^\d+\.\d+', line):
            return 2
        elif re.match(r'^\d+\.', line):
            return 1
        else:
            return 1  # Default to level 1
    
    def _is_list_item(self, line: str) -> bool:
        """Determine if a line is a list item."""
        list_patterns = [
            r'^[-*•]\s+',        # Bullet points
            r'^\d+\.\s+',        # Numbered lists
            r'^[a-zA-Z]\.\s+',   # Lettered lists
            r'^\(\d+\)\s+',      # Parenthetical numbers
        ]
        
        for pattern in list_patterns:
            if re.match(pattern, line):
                return True
        
        return False
    
    async def _perform_enhanced_analysis(self, structured_content: Dict[str, Any], 
                                        file_path: str) -> Optional[EnhancedTextAnalysis]:
        """Perform enhanced AI analysis with key concept extraction and confidence scoring."""
        try:
            # Get the appropriate model for analysis
            model_spec = self.model_manager.get_model_for_agent("text", "text")
            if not model_spec:
                raise Exception("No suitable model available for text analysis")
            
            raw_text = structured_content.get("raw_text", "")
            structure = structured_content.get("structure")
            
            # Perform key concept extraction
            key_concepts = await self._extract_key_concepts(raw_text, model_spec)
            
            # Generate educational metadata
            educational_metadata = await self._generate_educational_metadata(
                raw_text, structure, model_spec
            )
            
            # Calculate confidence scores
            confidence_scores = self._calculate_confidence_scores(
                raw_text, key_concepts, educational_metadata
            )
            
            # Collect processing metrics
            processing_metrics = {
                "model_used": model_spec.model_id,
                "text_length": len(raw_text),
                "concepts_extracted": len(key_concepts),
                "structure_elements": {
                    "headings": len(structure.headings) if structure else 0,
                    "paragraphs": len(structure.paragraphs) if structure else 0,
                    "lists": len(structure.lists) if structure else 0,
                    "tables": len(structure.tables) if structure else 0
                }
            }
            
            return EnhancedTextAnalysis(
                structure=structure or TextStructure([], [], [], [], {}),
                key_concepts=key_concepts,
                educational_metadata=educational_metadata,
                confidence_scores=confidence_scores,
                processing_metrics=processing_metrics
            )
            
        except Exception as e:
            print(f"⚠️ Error in enhanced analysis: {e}")
            # Try fallback analysis
            return await self._fallback_analysis(structured_content, file_path, str(e))
    
    async def _extract_key_concepts(self, text: str, model_spec) -> List[KeyConcept]:
        """Extract key concepts with confidence scoring using AI."""
        try:
            # Limit text for analysis to avoid token limits
            analysis_text = text[:3000] if len(text) > 3000 else text
            
            prompt = f"""
Analyze the following text and extract key concepts with detailed information:

Text: {analysis_text}

For each key concept, provide:
1. The concept name/term
2. A clear definition or explanation
3. Importance score (0.0 to 1.0) - how important this concept is to understanding the content
4. Confidence score (0.0 to 1.0) - how confident you are in this extraction
5. Context - where/how this concept appears in the text

Return ONLY a valid JSON array with this structure:
[
  {{
    "concept": "concept name",
    "definition": "clear definition",
    "importance": 0.8,
    "confidence": 0.9,
    "context": "brief context"
  }}
]

Extract 5-10 most important concepts. Focus on educational or technical terms that are central to understanding the content.
"""
            
            response = await self._invoke_model(model_spec, prompt, max_tokens=2000)
            
            # Parse the JSON response
            try:
                concepts_data = json.loads(response)
                key_concepts = []
                
                for concept_data in concepts_data:
                    if isinstance(concept_data, dict) and all(
                        key in concept_data for key in ['concept', 'definition', 'importance', 'confidence', 'context']
                    ):
                        key_concepts.append(KeyConcept(
                            concept=concept_data['concept'],
                            definition=concept_data['definition'],
                            importance=float(concept_data['importance']),
                            confidence=float(concept_data['confidence']),
                            context=concept_data['context']
                        ))
                
                return key_concepts
                
            except json.JSONDecodeError:
                print("⚠️ Failed to parse key concepts JSON, using fallback extraction")
                return self._fallback_concept_extraction(text)
            
        except Exception as e:
            print(f"⚠️ Error extracting key concepts: {e}")
            return self._fallback_concept_extraction(text)
    
    async def _generate_educational_metadata(self, text: str, structure, model_spec) -> Dict[str, Any]:
        """Generate educational metadata using AI analysis."""
        try:
            # Prepare structure summary
            structure_summary = ""
            if structure:
                structure_summary = f"""
Structure Analysis:
- Headings: {len(structure.headings)} found
- Paragraphs: {len(structure.paragraphs)}
- Lists: {len(structure.lists)}
- Tables: {len(structure.tables)}
"""
            
            analysis_text = text[:2500] if len(text) > 2500 else text
            
            prompt = f"""
Analyze this educational content and generate comprehensive metadata:

{structure_summary}

Content: {analysis_text}

Generate educational metadata in JSON format:
{{
  "learning_objectives": ["objective 1", "objective 2", ...],
  "key_topics": ["topic 1", "topic 2", ...],
  "difficulty_level": "beginner|intermediate|advanced",
  "estimated_reading_time": minutes_as_integer,
  "target_audience": "description of intended audience",
  "prerequisites": ["prerequisite 1", "prerequisite 2", ...],
  "bloom_taxonomy_levels": ["remember", "understand", "apply", "analyze", "evaluate", "create"],
  "content_type": "lecture|tutorial|reference|exercise|other",
  "summary": "2-3 sentence summary",
  "main_themes": ["theme 1", "theme 2", ...]
}}

Focus on educational value and learning outcomes. Be specific and actionable.
"""
            
            response = await self._invoke_model(model_spec, prompt, max_tokens=1500)
            
            try:
                metadata = json.loads(response)
                # Validate and clean the metadata
                return self._validate_educational_metadata(metadata)
                
            except json.JSONDecodeError:
                print("⚠️ Failed to parse educational metadata JSON")
                return self._generate_fallback_metadata(text)
            
        except Exception as e:
            print(f"⚠️ Error generating educational metadata: {e}")
            return self._generate_fallback_metadata(text)
    
    def _calculate_confidence_scores(self, text: str, key_concepts: List[KeyConcept], 
                                   educational_metadata: Dict[str, Any]) -> Dict[str, float]:
        """Calculate confidence scores for different aspects of the analysis."""
        scores = {}
        
        # Text quality score
        text_length = len(text)
        if text_length > 1000:
            scores["text_quality"] = 0.9
        elif text_length > 500:
            scores["text_quality"] = 0.7
        else:
            scores["text_quality"] = 0.5
        
        # Concept extraction confidence
        if key_concepts:
            avg_concept_confidence = sum(c.confidence for c in key_concepts) / len(key_concepts)
            scores["concept_extraction"] = avg_concept_confidence
        else:
            scores["concept_extraction"] = 0.3
        
        # Educational metadata confidence
        required_fields = ["learning_objectives", "difficulty_level", "summary"]
        metadata_completeness = sum(
            1 for field in required_fields 
            if field in educational_metadata and educational_metadata[field]
        ) / len(required_fields)
        scores["educational_metadata"] = metadata_completeness
        
        # Overall confidence
        scores["overall"] = sum(scores.values()) / len(scores)
        
        return scores
    
    async def _invoke_model(self, model_spec, prompt: str, max_tokens: int = 1000) -> str:
        """Invoke the specified model with fallback handling."""
        try:
            response = self.bedrock_client.invoke_model(
                modelId=model_spec.model_id,
                body=json.dumps({
                    "anthropic_version": "bedrock-2023-05-31",
                    "max_tokens": max_tokens,
                    "temperature": model_spec.temperature,
                    "messages": [{"role": "user", "content": prompt}]
                })
            )
            
            result = json.loads(response['body'].read())
            return result['content'][0]['text']
            
        except Exception as e:
            print(f"⚠️ Model invocation failed: {e}")
            # Try fallback model
            fallback_model = self.model_manager.get_fallback_model("text", model_spec.model_id)
            if fallback_model and fallback_model.model_id != model_spec.model_id:
                print(f"🔄 Trying fallback model: {fallback_model.model_id}")
                return await self._invoke_model(fallback_model, prompt, max_tokens)
            else:
                raise e
    
    def _get_current_model_id(self) -> str:
        """Get the current model ID being used."""
        model_spec = self.model_manager.get_model_for_agent("text")
        return model_spec.model_id if model_spec else "unknown"
    
    def _fallback_concept_extraction(self, text: str) -> List[KeyConcept]:
        """Fallback method for concept extraction using simple heuristics."""
        concepts = []
        
        # Simple keyword extraction based on capitalization and frequency
        words = re.findall(r'\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*\b', text)
        word_freq = {}
        
        for word in words:
            if len(word) > 3:  # Ignore short words
                word_freq[word] = word_freq.get(word, 0) + 1
        
        # Get top concepts by frequency
        sorted_concepts = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:5]
        
        for concept, freq in sorted_concepts:
            concepts.append(KeyConcept(
                concept=concept,
                definition=f"Key term appearing {freq} times in the text",
                importance=min(freq / 10, 1.0),
                confidence=0.6,  # Lower confidence for fallback method
                context="Extracted using frequency analysis"
            ))
        
        return concepts
    
    def _generate_fallback_metadata(self, text: str) -> Dict[str, Any]:
        """Generate basic metadata using simple heuristics."""
        word_count = len(text.split())
        estimated_reading_time = max(1, word_count // 200)  # Assume 200 words per minute
        
        # Simple difficulty assessment based on text characteristics
        avg_word_length = sum(len(word) for word in text.split()) / max(len(text.split()), 1)
        if avg_word_length > 6:
            difficulty = "advanced"
        elif avg_word_length > 4:
            difficulty = "intermediate"
        else:
            difficulty = "beginner"
        
        return {
            "learning_objectives": ["Understand the main concepts presented in the text"],
            "key_topics": ["General content analysis"],
            "difficulty_level": difficulty,
            "estimated_reading_time": estimated_reading_time,
            "target_audience": "General learners",
            "prerequisites": [],
            "bloom_taxonomy_levels": ["remember", "understand"],
            "content_type": "other",
            "summary": f"Text document with approximately {word_count} words covering various topics.",
            "main_themes": ["Content analysis"]
        }
    
    def _validate_educational_metadata(self, metadata: Dict[str, Any]) -> Dict[str, Any]:
        """Validate and clean educational metadata."""
        # Ensure required fields exist with defaults
        defaults = {
            "learning_objectives": [],
            "key_topics": [],
            "difficulty_level": "intermediate",
            "estimated_reading_time": 5,
            "target_audience": "General learners",
            "prerequisites": [],
            "bloom_taxonomy_levels": ["remember", "understand"],
            "content_type": "other",
            "summary": "Educational content analysis",
            "main_themes": []
        }
        
        for key, default_value in defaults.items():
            if key not in metadata or not metadata[key]:
                metadata[key] = default_value
        
        # Validate difficulty level
        valid_difficulties = ["beginner", "intermediate", "advanced"]
        if metadata["difficulty_level"] not in valid_difficulties:
            metadata["difficulty_level"] = "intermediate"
        
        # Ensure reading time is reasonable
        if not isinstance(metadata["estimated_reading_time"], int) or metadata["estimated_reading_time"] < 1:
            metadata["estimated_reading_time"] = 5
        
        return metadata
    
    async def _fallback_processing(self, file_path: str, error_msg: str) -> Dict[str, Any]:
        """Fallback processing when enhanced analysis fails."""
        try:
            # Resolve file path first
            resolved_path = self._resolve_file_path(file_path)
            
            # Try basic text extraction
            with open(resolved_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            return {
                "agent_type": "text",
                "file_path": file_path,
                "status": "completed_with_fallback",
                "content": {
                    "text": self._create_smart_preview(content, 8000),
                    "full_text": content,
                    "word_count": len(content.split()),
                    "char_count": len(content)
                },
                "enhanced_analysis": {},
                "metadata": {
                    "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0,
                    "file_type": Path(file_path).suffix.lower(),
                    "processed_by": "text_agent_fallback",
                    "error": error_msg,
                    "processing_version": "1.0_fallback"
                }
            }
            
        except Exception as fallback_error:
            return {
                "agent_type": "text",
                "file_path": file_path,
                "status": "error",
                "error": f"Primary error: {error_msg}, Fallback error: {str(fallback_error)}"
            }
    
    async def _fallback_analysis(self, structured_content: Dict[str, Any], 
                                file_path: str, error_msg: str) -> Optional[EnhancedTextAnalysis]:
        """Fallback analysis when enhanced analysis fails."""
        try:
            raw_text = structured_content.get("raw_text", "")
            
            # Use simple fallback methods
            key_concepts = self._fallback_concept_extraction(raw_text)
            educational_metadata = self._generate_fallback_metadata(raw_text)
            confidence_scores = {"overall": 0.5, "fallback_used": True}
            
            structure = structured_content.get("structure") or TextStructure([], [], [], [], {})
            
            processing_metrics = {
                "model_used": "fallback_analysis",
                "text_length": len(raw_text),
                "concepts_extracted": len(key_concepts),
                "fallback_reason": error_msg
            }
            
            return EnhancedTextAnalysis(
                structure=structure,
                key_concepts=key_concepts,
                educational_metadata=educational_metadata,
                confidence_scores=confidence_scores,
                processing_metrics=processing_metrics
            )
            
        except Exception as e:
            print(f"⚠️ Fallback analysis also failed: {e}")
            return None

    def _create_smart_preview(self, content: str, max_length: int = 8000) -> str:
        """Create a smart preview that preserves important content instead of simple truncation."""
        if len(content) <= max_length:
            return content
        
        # Split content into paragraphs
        paragraphs = content.split('\n\n')
        
        # Prioritize paragraphs that contain important information
        important_paragraphs = []
        regular_paragraphs = []
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            # Check if paragraph contains important keywords or patterns
            if (any(keyword in para.lower() for keyword in 
                   ['introduction', 'conclusion', 'summary', 'important', 'key', 'definition', 
                    'overview', 'objective', 'goal', 'purpose', 'main', 'primary', 'essential']) or
                para.startswith('#') or  # Headings
                ':' in para[:100] or  # Likely definitions or key points
                len(para.split()) < 50):  # Short paragraphs are often important
                important_paragraphs.append(para)
            else:
                regular_paragraphs.append(para)
        
        # Build preview starting with important paragraphs
        preview_parts = []
        current_length = 0
        
        # Add important paragraphs first
        for para in important_paragraphs:
            if current_length + len(para) > max_length:
                break
            preview_parts.append(para)
            current_length += len(para) + 2  # +2 for \n\n
        
        # Add regular paragraphs until we reach the limit
        for para in regular_paragraphs:
            if current_length + len(para) > max_length:
                break
            preview_parts.append(para)
            current_length += len(para) + 2
        
        # If we still have space, add a truncation indicator
        result = '\n\n'.join(preview_parts)
        if len(content) > len(result):
            result += f"\n\n[Content continues... {len(content) - len(result)} more characters]"
        
        return result


# Global instance
text_agent = TextAgent()